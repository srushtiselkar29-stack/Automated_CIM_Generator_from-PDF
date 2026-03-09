from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pathlib import Path


class PPTGenerator:

    def __init__(self):
        self.output_dir = "storage/generated_ppt"
        Path(self.output_dir).mkdir(parents=True, exist_ok=True)

    def create_presentation(self, title: str, slides: list) -> str:

        prs = Presentation()

        # -------------------------
        # Title Slide
        # -------------------------

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = "Confidential Information Memorandum"

        # -------------------------
        # Content Slides
        # -------------------------

        for slide_data in slides:

            visual = slide_data.get("visual", "bullets")

            if visual == "image":
                self._create_image_slide(prs, slide_data)

            elif visual == "chart":
                self._create_chart_slide(prs, slide_data)

            elif visual == "metrics":
                self._create_metric_slide(prs, slide_data)

            else:
                self._create_bullet_slide(prs, slide_data)

        ppt_path = f"{self.output_dir}/generated_presentation.pptx"
        prs.save(ppt_path)

        return ppt_path

    # ----------------------------------------------------
    # Bullet Slide
    # ----------------------------------------------------

    def _create_bullet_slide(self, prs, slide_data):

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        lines = slide_data["content"].split("\n")

        for i, line in enumerate(lines):

            if i == 0:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()

            p.text = line
            p.level = 0

    # ----------------------------------------------------
    # Image Slide
    # ----------------------------------------------------

    def _create_image_slide(self, prs, slide_data):

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]

        if "content" in slide_data:
            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(1.5),
                Inches(5),
                Inches(3)
            )

            tf = textbox.text_frame
            tf.text = slide_data["content"]

        if "image_path" in slide_data:

            slide.shapes.add_picture(
                slide_data["image_path"],
                Inches(6),
                Inches(1.5),
                width=Inches(3)
            )

    # ----------------------------------------------------
    # Chart Slide
    # ----------------------------------------------------

    def _create_chart_slide(self, prs, slide_data):

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]
        
        chart_type = slide_data.get("chart_type", "bar")

        if chart_type == "line":
            chart_style = XL_CHART_TYPE.LINE

        elif chart_type == "pie":
            chart_style = XL_CHART_TYPE.PIE

        elif chart_type == "area":
            chart_style = XL_CHART_TYPE.AREA

        else:
            chart_style = XL_CHART_TYPE.COLUMN_CLUSTERED

        

        chart_data = ChartData()

        chart_data.categories = slide_data.get(
            "categories",
            ["2021", "2022", "2023", "2024"]
        )

        chart_data.add_series(
            "Revenue",
            slide_data.get("values", [100, 120, 150, 180])
        )

        slide.shapes.add_chart(
            chart_style,
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(4),
            chart_data
        )

    # ----------------------------------------------------
    # Metric Slide
    # ----------------------------------------------------

    def _create_metric_slide(self, prs, slide_data):

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]

        metrics = slide_data.get("metrics", [])

        left = Inches(1)

        for metric in metrics:

            textbox = slide.shapes.add_textbox(
                left,
                Inches(2),
                Inches(3),
                Inches(2)
            )

            tf = textbox.text_frame

            p = tf.paragraphs[0]
            p.text = metric
            p.font.size = Pt(36)

            left += Inches(3)